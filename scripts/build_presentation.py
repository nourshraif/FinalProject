"""Generate the Vertex graduation-project PowerPoint deck.

Run:  python scripts/build_presentation.py
Output: Vertex_Presentation.pptx in the project root.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT.parent / ".cursor" / "projects" / "c-Users-nours-Desktop-FinalProject" / "assets"
DIAGRAM = ASSETS / "vertex_architecture.png"
OUT = ROOT / "Vertex_Presentation.pptx"

# Palette (dark navy / indigo)
BG = RGBColor(0x0B, 0x12, 0x24)
CARD = RGBColor(0x12, 0x1C, 0x33)
PRIMARY = RGBColor(0x8B, 0x8C, 0xF0)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
WHITE = RGBColor(0xF5, 0xF7, 0xFF)
MUTED = RGBColor(0xA9, 0xB2, 0xC9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    shape = slide.shapes.add_shape(
        1, 0, 0, SLIDE_W, SLIDE_H  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    shape.shadow.inherit = False
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    return shape


def add_text(slide, left, top, width, height, lines, *, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, bullet=False,
             line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = (f"•  {line}" if bullet else line)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return box


def add_badge(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(4.5), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run.font.name = "Segoe UI"
    return box


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(badge, title, bullets, notes, *, demo=False):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_badge(slide, badge)
    # accent bar
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.55), Inches(1.4), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
             title, size=34, bold=True, color=WHITE)
    add_text(slide, Inches(0.7), Inches(1.9), Inches(12), Inches(4.8),
             bullets, size=20, color=MUTED, bullet=True, line_spacing=1.3)
    if demo:
        add_text(slide, Inches(0.7), Inches(5.9), Inches(12), Inches(0.8),
                 "LIVE DEMO", size=22, bold=True, color=PRIMARY)
    add_notes(slide, notes)
    return slide


# ── Slide 1 — Title ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_text(s, Inches(0), Inches(2.3), SLIDE_W, Inches(1.2),
         "Vertex", size=66, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), SLIDE_W, Inches(0.7),
         "AI-Powered Job Matching Platform", size=26, color=PRIMARY,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.6), SLIDE_W, Inches(1.6),
         ["Final Year Graduation Project",
          "Presented by: [Name 1]  ·  [Name 2]  ·  [Name 3]",
          "Supervisor: [Name]  ·  [University]  ·  2026"],
         size=18, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.4)
add_notes(s, "Good morning. We're presenting Vertex, an AI-powered platform "
             "that matches job seekers to the right jobs. I'm [name], and I'll "
             "start with the problem we set out to solve.")

# ── Slide 2 — Problem (Presenter 1) ─────────────────────────────
content_slide(
    "PRESENTER 1  ·  CONTEXT",
    "The Problem",
    ["Job listings are scattered across dozens of separate job boards",
     "Traditional search is keyword-based — it misses qualified candidates",
     "Job seekers waste hours scrolling through irrelevant postings",
     "Companies struggle to reach the right applicants"],
    "Anyone job-hunting knows the pain: you check five different sites, and "
    "keyword search either floods you with irrelevant results or hides jobs "
    "you'd actually be great for.")

# ── Slide 3 — Solution (Presenter 1) ────────────────────────────
content_slide(
    "PRESENTER 1  ·  SOLUTION",
    "Our Solution",
    ["Vertex reads your CV and understands it by meaning, not keywords",
     "Aggregates jobs from 8 job boards into one place",
     "Recommends roles you're actually qualified for",
     "Companies can also post jobs directly on the platform"],
    "Vertex solves this by aggregating jobs into one place and using AI to "
    "semantically match your CV to roles you're genuinely qualified for. Now "
    "[Name 2] will walk you through the features.")

# ── Slide 4 — Features (Presenter 2) ────────────────────────────
content_slide(
    "PRESENTER 2  ·  FEATURES",
    "Core Features",
    ["CV upload + AI parsing — extracts skills & experience",
     "Semantic matching — ranks jobs by relevance to your profile",
     "Job search & filters — by board, location, date posted",
     "Company portal — register and post jobs",
     "AI chatbot assistant + user & admin dashboards"],
    "Users upload a CV and instantly get matched jobs. Companies get their own "
    "portal to post roles, and admins manage everything from a dashboard.")

# ── Slide 5 — Demo 1 (Presenter 2) ──────────────────────────────
content_slide(
    "PRESENTER 2  ·  DEMO",
    "Live Demo #1 — CV Matching",
    ["Upload a CV → see semantic job matches",
     "Show a top match and why it matched the profile"],
    "Rather than just tell you, let me show it live. (Upload CV, show ranked "
    "matches on the deployed site.)",
    demo=True)

# ── Slide 6 — Technologies (Presenter 2) ────────────────────────
content_slide(
    "PRESENTER 2  ·  TECHNOLOGIES",
    "Technologies Used",
    ["Frontend: Next.js, React, TypeScript, Tailwind CSS",
     "Backend: Python, FastAPI",
     "Database: PostgreSQL + pgvector (vector similarity)",
     "AI: text embeddings for semantic matching",
     "Auth: Google OAuth 2.0 + email/password",
     "Deployment: Docker, Nginx, HTTPS, DigitalOcean VPS"],
    "Our stack: a Next.js frontend, a FastAPI backend, and Postgres with "
    "pgvector for the AI matching. Everything runs in Docker on a live server. "
    "[Name 3] will explain how it all fits together and the challenges we hit.")

# ── Slide 7 — Architecture (Presenter 3) ────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_badge(s, "PRESENTER 3  ·  ARCHITECTURE")
add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
         "How It Works", size=34, bold=True, color=WHITE)
if DIAGRAM.exists():
    img_w = Inches(10.5)
    left = int((SLIDE_W - img_w) / 2)
    s.shapes.add_picture(str(DIAGRAM), left, Inches(1.9), width=img_w)
else:
    add_text(s, Inches(0.7), Inches(2), Inches(12), Inches(3),
             ["[Insert vertex_architecture.png here]"], size=20, color=MUTED)
add_notes(s, "Here's the architecture. A nightly pipeline pulls jobs from 8 "
             "boards, cleans and deduplicates them, and stores vector "
             "embeddings. When a user uploads a CV, we embed it too and use "
             "pgvector to find the closest matches.")

# ── Slide 8 — Challenges (Presenter 3) ──────────────────────────
content_slide(
    "PRESENTER 3  ·  CHALLENGES",
    "Main Challenges & Solutions",
    ["Match quality: keyword search failed → vector embeddings + cosine similarity",
     "Broken scrapers: sites blocked us / used JS → switched to official JSON APIs",
     "Duplicate & stale jobs → deduplication + 30-day TTL cleanup",
     "Deployment: localhost API, DB init, HTTPS conflicts → same-origin API, "
     "init on startup, Let's Encrypt in Nginx"],
    "Our biggest challenges were AI match quality — solved with vector "
    "embeddings — and unreliable scrapers, which we fixed by moving to official "
    "APIs. Deployment also took real problem-solving to get HTTPS and the "
    "database working in production.")

# ── Slide 9 — Demo 2 (Presenter 3) ──────────────────────────────
content_slide(
    "PRESENTER 3  ·  DEMO",
    "Live Demo #2 — Deployed Platform",
    ["Open the live site: vertex-grad.duckdns.org",
     "Show aggregated real jobs + filter by job board"],
    "And this is all live in production. (Open the site, filter jobs by board.)",
    demo=True)

# ── Slide 10 — Conclusion (Presenter 3) ─────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_badge(s, "PRESENTER 3  ·  CONCLUSION")
add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
         "Conclusion", size=34, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(1.9), Inches(12), Inches(3),
         ["Vertex: aggregation + AI matching in one deployed platform",
          "Learned: full-stack development, applied AI, real-world deployment",
          "Future work: more job boards, mobile app, richer analytics"],
         size=20, color=MUTED, bullet=True, line_spacing=1.35)
add_text(s, Inches(0), Inches(5.4), SLIDE_W, Inches(1),
         "Thank you — Questions?", size=30, bold=True, color=PRIMARY,
         align=PP_ALIGN.CENTER)
add_notes(s, "To wrap up, Vertex brings job aggregation and AI matching "
             "together in a fully deployed product. This project taught us "
             "full-stack development, applied AI, and real-world deployment. "
             "Thank you — we're happy to take questions.")

prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides._sldIdLst)}")
